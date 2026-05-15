"""Generate a single-slide PPT visualisation of PM v3 architecture.

Output: C:/Users/26636/Desktop/MyCrew-PM-v3-Architecture.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt


# ── Theme tokens ─────────────────────────────────────────────────────

INK = RGBColor(0x1F, 0x29, 0x37)          # 深墨灰
INK_SOFT = RGBColor(0x4B, 0x55, 0x63)
INK_MUTED = RGBColor(0x6B, 0x72, 0x80)
INK_GHOST = RGBColor(0x9C, 0xA3, 0xAF)

BRAND = RGBColor(0x0C, 0x8C, 0xE9)        # 品牌蓝
BRAND_SOFT = RGBColor(0xE3, 0xF2, 0xFD)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_SOFT = RGBColor(0xD1, 0xFA, 0xE5)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_SOFT = RGBColor(0xFE, 0xF3, 0xC7)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_SOFT = RGBColor(0xEE, 0xE6, 0xFF)
SURFACE = RGBColor(0xFA, 0xFB, 0xFC)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)


# ── Geometry ─────────────────────────────────────────────────────────

SLIDE_W = Emu(int(13.333 * 914400))   # 13.33" widescreen
SLIDE_H = Emu(int(7.5 * 914400))


def inches(v: float) -> Emu:
    return Emu(int(v * 914400))


# ── Helpers ──────────────────────────────────────────────────────────

def add_text(slide, x, y, w, h, text, *, size=10, bold=False, color=INK,
              align_center=False, italic=False, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        if align_center:
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
    return box


def add_rect(slide, x, y, w, h, *, fill=CARD, line=BORDER, line_w=0.75, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    if radius and hasattr(shape, "adjustments"):
        try:
            shape.adjustments[0] = 0.10
        except Exception:
            pass
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=INK_MUTED, weight=1.5):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # Add arrow head
    line_elem = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tail = line_elem.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return conn


# ── Build slide ──────────────────────────────────────────────────────


def build_slide() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H,
    )
    bg.fill.solid(); bg.fill.fore_color.rgb = SURFACE
    bg.line.fill.background()
    bg.shadow.inherit = False

    # ── Title bar ───────────────────────────────────────────────────
    add_text(slide, inches(0.5), inches(0.25), inches(12.3), inches(0.5),
             "MyCrew PM v3 — 5-Phase Crew Architecture",
             size=22, bold=True, color=INK)
    add_text(slide, inches(0.5), inches(0.7), inches(12.3), inches(0.3),
             "create_new sub-agent redesigned: strict Pydantic contracts per phase + in-memory cache + breakpoint resume + transparent debug log",
             size=10, italic=True, color=INK_MUTED)

    # ── Pre-pipeline (Router) ───────────────────────────────────────
    pre_y = inches(1.25)
    pre_h = inches(0.5)
    # User → Router → intent_classifier
    boxes = [
        ("用户消息", BRAND_SOFT, BRAND),
        ("pre_filter\n(正则)", CARD, BORDER),
        ("compliance_gate\n(cheap)", CARD, BORDER),
        ("intent_classifier\n(cheap)", CARD, BORDER),
        ("create_new\n(thin entry)", PURPLE_SOFT, PURPLE),
    ]
    box_w = inches(2.05)
    gap = inches(0.15)
    start_x = inches(0.5)
    for i, (label, fill, line) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, pre_y, box_w, pre_h, fill=fill, line=line, line_w=1.0)
        add_text(slide, x, pre_y + inches(0.05), box_w, pre_h - inches(0.1),
                 label, size=10, bold=True, color=INK_SOFT, align_center=True)
        if i < len(boxes) - 1:
            ax = x + box_w
            ay = pre_y + pre_h // 2
            add_arrow(slide, ax, ay, ax + gap, ay, color=INK_GHOST, weight=1.2)

    # ── 5-phase crew (the core) ─────────────────────────────────────
    crew_y = inches(2.15)
    crew_h = inches(2.6)
    # Frame
    add_rect(slide, inches(0.5), crew_y, inches(12.3), crew_h,
             fill=CARD, line=BRAND, line_w=1.5)
    add_text(slide, inches(0.65), crew_y + inches(0.05),
             inches(12), inches(0.35),
             "_planner_orchestrator.run_crew  ·  5 sequential CrewAI kickoffs  ·  each emits pm.log WS events",
             size=10, bold=True, color=BRAND)

    # Phase 0 (full width)
    p0_y = crew_y + inches(0.45)
    p0_h = inches(0.55)
    add_rect(slide, inches(0.8), p0_y, inches(11.7), p0_h,
             fill=BRAND_SOFT, line=BRAND, line_w=0.75)
    add_text(slide, inches(0.95), p0_y + inches(0.05),
             inches(11.5), p0_h - inches(0.1),
             "Phase 0 · 完整度判定   cheap LLM, t=0, max_tokens=10   →   ONELINE | PRD   (PRD 跳过 Phase 1)",
             size=10, bold=True, color=BRAND)

    # Phase 1-5 boxes
    phases = [
        ("Phase 1\n游戏主策划", "submit_concept", "ConceptDoc", "pro, t=0.7"),
        ("Phase 2\n系统策划", "submit_atomic_tasks", "AtomicTask[]", "pro, t=0.5"),
        ("Phase 3\n审核策划", "submit_reviewed_tasks", "ReviewedTask[]", "pro, t=0.2"),
        ("Phase 4\n项目管理", "submit_pathed_tasks", "PathedTask[]\n+ setup task", "pro, t=0.3"),
        ("Phase 5\nAgent 指挥员", "submit_assignments", "Assignment[]", "pro, t=0.2"),
    ]
    phase_y = p0_y + p0_h + inches(0.15)
    phase_h = inches(1.4)
    phase_w = inches(2.3)
    phase_gap = inches(0.05)
    phase_start = inches(0.8)
    colors = [
        (PURPLE_SOFT, PURPLE),
        (PURPLE_SOFT, PURPLE),
        (GREEN_SOFT, GREEN),
        (AMBER_SOFT, AMBER),
        (BRAND_SOFT, BRAND),
    ]
    for i, ((role, tool, schema, llm_cfg), (fill, line)) in enumerate(zip(phases, colors)):
        x = phase_start + i * (phase_w + phase_gap)
        add_rect(slide, x, phase_y, phase_w, phase_h, fill=fill, line=line, line_w=1.0)
        # Header (role)
        add_text(slide, x + inches(0.05), phase_y + inches(0.05),
                 phase_w - inches(0.1), inches(0.45), role,
                 size=10, bold=True, color=INK, align_center=True)
        # Tool name
        add_text(slide, x + inches(0.05), phase_y + inches(0.5),
                 phase_w - inches(0.1), inches(0.25),
                 f"⚙ {tool}",
                 size=8, color=INK_SOFT, align_center=True, font="Consolas")
        # Output schema
        add_text(slide, x + inches(0.05), phase_y + inches(0.72),
                 phase_w - inches(0.1), inches(0.45),
                 f"→ {schema}",
                 size=8, color=INK_SOFT, align_center=True, font="Consolas")
        # LLM config
        add_text(slide, x + inches(0.05), phase_y + inches(1.15),
                 phase_w - inches(0.1), inches(0.2),
                 llm_cfg, size=7, italic=True,
                 color=INK_GHOST, align_center=True)
        if i < len(phases) - 1:
            ax = x + phase_w
            ay = phase_y + phase_h // 2
            add_arrow(slide, ax, ay, ax + phase_gap, ay,
                       color=INK_MUTED, weight=1.5)

    # ── Cache + Persist row ─────────────────────────────────────────
    cache_y = crew_y + crew_h + inches(0.2)
    cache_h = inches(1.1)

    # Cache box
    add_rect(slide, inches(0.5), cache_y, inches(5.5), cache_h,
             fill=GREEN_SOFT, line=GREEN, line_w=1.0)
    add_text(slide, inches(0.65), cache_y + inches(0.05),
             inches(5.2), inches(0.3),
             "📦 planner_cache_svc  ·  in-memory dict",
             size=11, bold=True, color=GREEN)
    add_text(slide, inches(0.65), cache_y + inches(0.35),
             inches(5.2), cache_h - inches(0.4),
             "_sessions[session_id] {\n"
             "   status, phase_outputs, debug_log, draft_blueprint,\n"
             "   cancel_requested, pm_task\n"
             "}\n"
             "全程 RAM；关程序自动消亡；新会话清空",
             size=8, color=INK_SOFT, font="Consolas")

    # Save button arrow
    arrow_y = cache_y + cache_h // 2
    add_arrow(slide, inches(6.0), arrow_y, inches(6.8), arrow_y,
               color=BRAND, weight=2.0)
    add_text(slide, inches(6.0), arrow_y - inches(0.28),
             inches(0.8), inches(0.22),
             "保存项目", size=9, bold=True,
             color=BRAND, align_center=True)

    # Persist box
    add_rect(slide, inches(6.8), cache_y, inches(6.0), cache_h,
             fill=BRAND_SOFT, line=BRAND, line_w=1.0)
    add_text(slide, inches(6.95), cache_y + inches(0.05),
             inches(5.7), inches(0.3),
             "💾 planner_persist_svc.save_draft_as_project",
             size=11, bold=True, color=BRAND)
    add_text(slide, inches(6.95), cache_y + inches(0.35),
             inches(5.7), cache_h - inches(0.4),
             "1. project_svc.create_project_with_tasks  →  DB\n"
             "2. blueprint_writer.write_blueprint_to_disk  →  .mycrew/\n"
             "3. dump _planner_trace.json  (forensics)\n"
             "4. session.project_id ← project_id ; broadcast workflow_created",
             size=8, color=INK_SOFT, font="Consolas")

    # ── Failure path / breakpoint resume ────────────────────────────
    fail_y = cache_y + cache_h + inches(0.15)
    fail_h = inches(0.55)
    add_rect(slide, inches(0.5), fail_y, inches(12.3), fail_h,
             fill=AMBER_SOFT, line=AMBER, line_w=0.75)
    add_text(slide, inches(0.65), fail_y + inches(0.05),
             inches(12), fail_h - inches(0.1),
             "⚠️  失败时 3 层防御：①CrewAI self-correct (max_iter=5) → ②焦点修复 1 次 → "
             "③status=failed，前端「从断点重来」按钮 POST /pm/restart → "
             "上游 phase 产物从 cache 复用，仅重跑挂掉的 phase 起",
             size=9, bold=True, color=AMBER, align_center=False)

    # ── Bottom legend / key decisions ───────────────────────────────
    legend_y = fail_y + fail_h + inches(0.1)
    legend_h = inches(0.55)

    legend_items = [
        ("🔒 严格 Pydantic args_schema",
         "每 phase 一个 submit_xxx 工具，LLM 看到精确签名 → 一次过率拉高"),
        ("💾 Cache-First 持久化",
         "草稿全程 in-mem；用户点「保存项目」才入 DB + .mycrew/"),
        ("🪟 Drawer mount-keep",
         "关 drawer 不 unmount，所有 hooks 订阅活着；跨页面切换不打断"),
        ("📡 透明可观测",
         "pm.log WS 事件流；PMDebugLog 折叠式 phase 节；保存时 dump trace.json"),
    ]
    legend_w = (SLIDE_W - inches(1.0) - inches(0.3) * 3) // 4
    for i, (head, body) in enumerate(legend_items):
        x = inches(0.5) + i * (legend_w + inches(0.1))
        add_rect(slide, x, legend_y, legend_w, legend_h,
                 fill=CARD, line=BORDER, line_w=0.5)
        add_text(slide, x + inches(0.1), legend_y + inches(0.05),
                 legend_w - inches(0.2), inches(0.25),
                 head, size=9, bold=True, color=INK_SOFT)
        add_text(slide, x + inches(0.1), legend_y + inches(0.27),
                 legend_w - inches(0.2), inches(0.25),
                 body, size=7, color=INK_MUTED)

    return prs


def main() -> None:
    desktop = Path(r"C:/Users/26636/Desktop")
    out = desktop / "MyCrew-PM-v3-Architecture.pptx"
    prs = build_slide()
    prs.save(out)
    print(f"OK: {out}")


if __name__ == "__main__":
    main()
